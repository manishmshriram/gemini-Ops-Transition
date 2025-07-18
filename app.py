import streamlit as st
import time
from selenium import webdriver
from selenium.webdriver.edge.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.edge.options import Options
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

st.title("Gemini Ops Transition Reporting Generator")

def add_image_to_slide(slide, image_path, left_cm, top_cm, width_cm, height_cm):
    slide.shapes.add_picture(
        image_path,
        left=Cm(left_cm),
        top=Cm(top_cm),
        width=Cm(width_cm),
        height=Cm(height_cm),
    )

def add_rectangle_to_slide(slide, left_cm, top_cm, width_cm, height_cm):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(left_cm),
        Cm(top_cm),
        Cm(width_cm),
        Cm(height_cm)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.shadow.visible = False

def inject_confirmation_modal(driver):
    js_script = """
    if (!document.getElementById('chatgpt_ready_modal')) {
        let modal = document.createElement('div');
        modal.id = 'chatgpt_ready_modal';
        modal.style.position = 'fixed';
        modal.style.top = '0';
        modal.style.left = '0';
        modal.style.width = '100%';
        modal.style.height = '100%';
        modal.style.backgroundColor = 'rgba(0,0,0,0.5)';
        modal.style.display = 'flex';
        modal.style.alignItems = 'center';
        modal.style.justifyContent = 'center';
        modal.style.zIndex = '9999';

        let inner = document.createElement('div');
        inner.style.background = 'white';
        inner.style.padding = '20px';
        inner.style.borderRadius = '8px';
        inner.style.textAlign = 'center';
        inner.innerHTML = `
            <p style='font-size:18px;'>Is the webpage ready to run the script?</p>
            <button id='chatgpt_yes_button' style='margin:10px;padding:10px;'>Yes</button>
            <button id='chatgpt_no_button' style='margin:10px;padding:10px;'>No</button>
        `;

        modal.appendChild(inner);
        document.body.appendChild(modal);

        document.getElementById('chatgpt_yes_button').onclick = function() {
            modal.setAttribute('data-response', 'yes');
            modal.style.display = 'none';
        }
        document.getElementById('chatgpt_no_button').onclick = function() {
            modal.setAttribute('data-response', 'no');
            modal.style.display = 'none';
        }
    } else {
        let modal = document.getElementById('chatgpt_ready_modal');
        modal.setAttribute('data-response', '');
        modal.style.display = 'flex';
    }
    """
    driver.execute_script(js_script)

def wait_for_user_confirmation(driver):
    inject_confirmation_modal(driver)
    while True:
        response = driver.execute_script("""
            let modal = document.getElementById('chatgpt_ready_modal');
            return modal ? modal.getAttribute('data-response') : null;
        """)
        if response == 'yes':
            print("✅ User confirmed to proceed.")
            break
        elif response == 'no':
            print("⏳ User requested delay, waiting 60 seconds before re-asking...")
            time.sleep(60)
            inject_confirmation_modal(driver)
        else:
            time.sleep(1)

if st.button("Run Scraper and Generate PPT"):
    st.info("Launching Edge browser and scraping... Please do not close the app.")

    options = Options()
    options.add_argument("--force-device-scale-factor=1.25")
    options.add_argument("--start-maximized")
    options.add_argument("--headless")
    options.add_argument("--disable-gpu")

    service = Service("msedgedriver.exe")
    driver = webdriver.Edge(service=service, options=options)
    driver_2 = webdriver.Edge(service=service, options=options)

    try:
        # Your scraping and PPT generation logic here...
        # Indent everything under try exactly as shown previously, 
        # fixing all mixed spaces and tabs into 4 spaces per level

        pass  # Placeholder to keep this example focused on indentation setup

    except Exception as e:
        st.error(f"An error occurred: {e}")

    finally:
        driver.quit()
        driver_2.quit()
        st.info("Browser sessions closed.")
