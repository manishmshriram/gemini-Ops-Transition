# Streamlit conversion of your Jupyter Selenium + PPT generator
# Ready for local deployment with fixed URL and no user inputs

import streamlit as st
import time
from selenium import webdriver
from selenium.webdriver.edge.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.action_chains import ActionChains
from selenium.webdriver.edge.options import Options
from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.enum.shapes import MSO_SHAPE
import os

st.title("Gemini Ops Transition Reporting Generator")
def add_image_to_slide(slide, image_path, left_cm, top_cm, width_cm, height_cm):
    """
    Add an image to a slide at a specific position and size.
    
    Parameters:
    - slide: The slide object where the image will be added.
    - image_path: Path to the image file.
    - left_cm: Left position in cm.
    - top_cm: Top position in cm.
    - width_cm: Width of the image in cm.
    - height_cm: Height of the image in cm.
    """
    # Convert cm to PowerPoint's internal units and add the image
    slide.shapes.add_picture(
        image_path,
        left=Cm(left_cm),
        top=Cm(top_cm),
        width=Cm(width_cm),
        height=Cm(height_cm),
    )
def add_rectangle_to_slide(slide, left_cm, top_cm, width_cm, height_cm):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(left_cm),
        Cm(top_cm),
        Cm(width_cm),
        Cm(height_cm)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.shadow.visible = False

def inject_confirmation_modal(driver):
    js_script = """
    if (!document.getElementById('chatgpt_ready_modal')) {
        let modal = document.createElement('div');
        modal.id = 'chatgpt_ready_modal';
        modal.style.position = 'fixed';
        modal.style.top = '0';
        modal.style.left = '0';
        modal.style.width = '100%';
        modal.style.height = '100%';
        modal.style.backgroundColor = 'rgba(0,0,0,0.5)';
        modal.style.display = 'flex';
        modal.style.alignItems = 'center';
        modal.style.justifyContent = 'center';
        modal.style.zIndex = '9999';

        let inner = document.createElement('div');
        inner.style.background = 'white';
        inner.style.padding = '20px';
        inner.style.borderRadius = '8px';
        inner.style.textAlign = 'center';
        inner.innerHTML = `
            <p style='font-size:18px;'>Is the webpage ready to run the script?</p>
            <button id='chatgpt_yes_button' style='margin:10px;padding:10px;'>Yes</button>
            <button id='chatgpt_no_button' style='margin:10px;padding:10px;'>No</button>
        `;

        modal.appendChild(inner);
        document.body.appendChild(modal);

        document.getElementById('chatgpt_yes_button').onclick = function() {
            modal.setAttribute('data-response', 'yes');
            modal.style.display = 'none';
        }
        document.getElementById('chatgpt_no_button').onclick = function() {
            modal.setAttribute('data-response', 'no');
            modal.style.display = 'none';
        }
    } else {
        // Reset for re-showing modal
        let modal = document.getElementById('chatgpt_ready_modal');
        modal.setAttribute('data-response', '');
        modal.style.display = 'flex';
    }
    """
    driver.execute_script(js_script)
def wait_for_user_confirmation(driver):
    inject_confirmation_modal(driver)
    while True:
        response = driver.execute_script("""
            let modal = document.getElementById('chatgpt_ready_modal');
            return modal ? modal.getAttribute('data-response') : null;
        """)
        if response == 'yes':
            print("✅ User confirmed to proceed.")
            break
        elif response == 'no':
            print("⏳ User requested delay, waiting 60 seconds before re-asking...")
            time.sleep(60)
            inject_confirmation_modal(driver)  # re-show modal
        else:
            time.sleep(1)

if st.button("Run Scraper and Generate PPT"):
    st.info("Launching Edge browser and scraping... Please do not close the app.")

    options = Options()
    options.add_argument("--force-device-scale-factor=1.25")
    options.add_argument("--start-maximized")
    options.add_argument("--headless")
    options.add_argument("--disable-gpu")

    # Path to your msedgedriver.exe in the same folder or provide full path
    service = Service("msedgedriver.exe")

    driver = webdriver.Edge(service=service, options=options)
    driver_2 = webdriver.Edge(service=service, options=options)
	
    try:
        # === YOUR FIXED URL LOGIC ===
        driver.get("https://hlag.sharepoint.com/sites/QlikSense/SitePages/Gemini-Transition-Reporting.aspx?web=1#operational-kpis")
        driver_2.get("https://qliksense.hlag.com/sense/app/6c33a525-1a3b-4d63-9d6e-d80c235910a7/sheet/58c38c9b-850a-433e-a671-e142bf439f59/state/analysis")
	driver.maximize_window()
	driver_2.maximize_window()
        # === Retain your sleep times exactly as provided ===
        time.sleep(60)  # example sleep, replace with your exact timing
        
        #from here
	 try:
        	WebDriverWait(driver_2, 60).until(EC.presence_of_element_located((By.XPATH, '/html/body/div[4]/div[4]/div/div/div[2]/div/div/div[3]/div/div/div[3]/div[2]/div[10]')))
        	print("Sheet_2 element is fully loaded!")
    	except TimeoutException:
        	print("Sheet_2 load took too long.")
    	time.sleep(120)
   	 wait = WebDriverWait(driver, 120)
    	try:
        	iframe = wait.until(EC.presence_of_element_located((By.XPATH, "/html/body/div[1]/div[2]/div[2]/div/div/div[3]/section/article/div/div/div/div/div[1]/div/div/div/div/div/div/div[20]/div/div/div/div/div[2]/div/div/div/div[2]/div/div/div/div/iframe")))
        	print("First Frame element is fully loaded!")
    	except TimeoutException:
        	print("First Frame took too long.")
    	wait_for_user_confirmation(driver)
    	port_1 = driver_2.find_element(By.XPATH, "/html/body/div[4]/div[4]/div/div/div[2]/div/div/div[3]/div/div/div[3]/div[2]/div[11]/div[1]/div/div[1]/article/div[1]/div/div/div/div/div/div[2]/div[2]/vzl-pivot-main-content")
    	port_2 = driver_2.find_element(By.XPATH, "/html/body/div[4]/div[4]/div/div/div[2]/div/div/div[3]/div/div/div[3]/div[2]/div[19]/div[1]/div/div[1]/article/div[1]/div/div/div/div/div/div[2]/div[2]/vzl-pivot-main-content")

    	time.sleep(5)
    
    	port_1.screenshot("port_1.png")
   	port_2.screenshot("port_2.png")
    	driver.switch_to.frame(iframe)
    

    	time.sleep(5)

    
    	image_1_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[8]/div/article/div[1]/div/div/div/div[5]/div[2]/div[1]/div/article/div[1]/div/div/div")
    	image_2_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[5]/div/article/div[1]/div/div/div")
    	image_3_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[4]/div/article/div[1]/div/div/div")
    	image_4_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[3]/div/article/div[1]/div/div/div")
    	image_5_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[8]/div/article/div[1]/div/div/div/div[5]/div[2]/div[1]/div/article/div[2]/footer")
    	image_6_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[5]/div/article/div[2]/footer")
    	image_7_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[4]/div/article/div[2]/footer")
    	image_8_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[3]/div/article/div[2]/footer")

    	image_1_element.screenshot("chart_1.png")
    	image_2_element.screenshot("chart_2.png")
    	image_3_element.screenshot("chart_3.png")
    	image_4_element.screenshot("chart_4.png")
    	image_5_element.screenshot("fotter_1.png")
    	image_6_element.screenshot("fotter_2.png")
    	image_7_element.screenshot("fotter_3.png")
    	image_8_element.screenshot("fotter_4.png")

    	time.sleep(2)
    
    	driver.switch_to.default_content()

    	time.sleep(2)
    
    	try:
        	iframe_2 = wait.until(EC.presence_of_element_located((By.XPATH, "/html/body/div[1]/div[2]/div[2]/div/div/div[3]/section/article/div/div/div/div/div[1]/div/div/div/div/div/div/div[20]/div/div/div/div/div[3]/div/div/div/div[2]/div/div/div/div/iframe")))
        	print("Second Frame element is fully loaded!")
    	except TimeoutException:
        	print("Second Frame took too long.")


    	driver.switch_to.frame(iframe_2)

    	time.sleep(5)

    	image_9_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[1]/div/article/div[1]/div")
    	image_10_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[1]/div/article/div[2]/footer")
    	image_9_element.screenshot("chart_5.png")
    	image_10_element.screenshot("fotter_5.png")

    	time.sleep(3)

    	driver.switch_to.default_content()

    	time.sleep(2)
    
    	try:
        	iframe_3 = wait.until(EC.presence_of_element_located((By.XPATH, "/html/body/div[1]/div[2]/div[2]/div/div/div[3]/section/article/div/div/div/div/div[1]/div/div/div/div/div/div/div[22]/div/div/div/div/div[2]/div/div/div/div[2]/div/div/div/div/iframe")))
        	print("Third Frame element is fully loaded!")
    	except TimeoutException:
        	print("Third Frame took too long.")


    	driver.switch_to.frame(iframe_3)

    	time.sleep(5)

    	image_11_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[1]/div/article/div[1]/div/div/div")
    	image_12_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[1]/div/article/div[2]/footer")
    	image_13_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[2]/div/article/div[2]/footer")
    	image_11_element.screenshot("chart_6.png")
    	image_12_element.screenshot("fotter_6.png")
    	image_13_element.screenshot("table_1_fotter.png")
    	element_c = driver.find_element(By.XPATH,"/html/body/div[4]/div/div[2]/div/article/div/div[2]/div/article/div[1]")
    	element_b = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[2]/div/article/div[1]/div/div/div/div[3]")
    	total_height_b = driver.execute_script("return arguments[0].scrollHeight", element_b)
    	viewport_height_b = driver.execute_script("return arguments[0].clientHeight", element_b)
    	# Element B TOP
    	driver.execute_script("arguments[0].scrollTop = 0", element_b)
    	time.sleep(0.5)
    	element_c.screenshot("table_1_1.png")
    	# Element B BOTTOM
    	driver.execute_script("arguments[0].scrollTop = arguments[0].scrollHeight", element_b)
    	time.sleep(0.5)
    	element_c.screenshot("table_1_2.png")
    	time.sleep(2)
       
   	driver.switch_to.default_content()

    	time.sleep(2)
    
    	try:
        	iframe_4 = wait.until(EC.presence_of_element_located((By.XPATH, "/html/body/div[1]/div[2]/div[2]/div/div/div[3]/section/article/div/div/div/div/div[1]/div/div/div/div/div/div/div[22]/div/div/div/div/div[3]/div/div/div/div[2]/div/div/div/div/iframe")))
        	print("Fourth Frame element is fully loaded!")
    	except TimeoutException:
        	print("Fourth Frame took too long.")


    	driver.switch_to.frame(iframe_4)
    	
    	time.sleep(5)

    	image_14_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[1]/div/article/div[1]/div/div/div")
    	image_15_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[1]/div/article/div[2]/footer")
    	image_16_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[2]/div/article/div[2]/footer")
    	image_14_element.screenshot("chart_7.png")
    	image_15_element.screenshot("fotter_7.png")
    	image_16_element.screenshot("table_2_fotter.png")
    	element_c = driver.find_element(By.XPATH,"/html/body/div[4]/div/div[2]/div/article/div/div[2]/div/article/div[1]")
    	element_b = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[2]/div/article/div[1]/div/div/div/div[3]")
    	total_height_b = driver.execute_script("return arguments[0].scrollHeight", element_b)
    	viewport_height_b = driver.execute_script("return arguments[0].clientHeight", element_b)
    	driver.execute_script("arguments[0].scrollTop = 0", element_b)
    	time.sleep(0.5)
    	element_c.screenshot("table_2_1.png")
    	# Element B BOTTOM
    	driver.execute_script("arguments[0].scrollTop = arguments[0].scrollHeight", element_b)
   	time.sleep(0.5)
    	element_c.screenshot("table_2_2.png")
    	time.sleep(2)
    	time.sleep(3)

    	driver.switch_to.default_content()

    	time.sleep(2)
    
    	try:
        	iframe_5 = wait.until(EC.presence_of_element_located((By.XPATH, "/html/body/div[1]/div[2]/div[2]/div/div/div[3]/section/article/div/div/div/div/div[1]/div/div/div/div/div/div/div[22]/div/div/div/div/div[4]/div/div/div/div[2]/div/div/div/div/iframe")))
        	print("Fifth Frame element is fully loaded!")
    	except TimeoutException:
        	print("Fifth Frame took too long.")


    	driver.switch_to.frame(iframe_5)

    	time.sleep(5)

    	image_17_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[1]/div/article/div[1]/div/div/div")
    	image_18_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[1]/div/article/div[2]/footer")
    	image_19_element = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[2]/div/article/div[2]/footer")
    	image_17_element.screenshot("chart_8.png")
    	image_18_element.screenshot("fotter_8.png")
    	image_19_element.screenshot("table_3_fotter.png")
    	element_c = driver.find_element(By.XPATH,"/html/body/div[4]/div/div[2]/div/article/div/div[2]/div/article/div[1]")
    	element_b = driver.find_element(By.XPATH, "/html/body/div[4]/div/div[2]/div/article/div/div[2]/div/article/div[1]/div/div/div/div[3]")
    	total_height_b = driver.execute_script("return arguments[0].scrollHeight", element_b)
    	viewport_height_b = driver.execute_script("return arguments[0].clientHeight", element_b)
    	driver.execute_script("arguments[0].scrollTop = 0", element_b)
    	time.sleep(0.5)
    	element_c.screenshot("table_3_1.png")
    	# Element B BOTTOM
    	driver.execute_script("arguments[0].scrollTop = arguments[0].scrollHeight", element_b)
    	time.sleep(0.5)
    	element_c.screenshot("table_3_2.png")

    	time.sleep(3)

    	driver.switch_to.default_content()
        
        	# === PPT Generation ===
        ppt = Presentation("Template.pptx")  # Replace with your template file path
    # Specify image details
    	image_path_1 = "chart_1.png" #Share Of Voyages Started On Time
    	image_path_2 = "chart_2.png" #All Gemini Port Calls | HL vs. Partner | Proforma SR [+/- 1 Day]
    	image_path_3 = "chart_3.png" #Global | Proforma SR [+/- 1 Day]
    	image_path_4 = "chart_4.png" #No. Of Gemini Port Calls | HL vs. Partner
    	image_path_5 = "chart_5.png" #Gemini By Trade | Proforma Schedule Reliability [+/- 1 Day] for Mainliner and Shuttles
    	image_path_6 = "chart_6.png" #Hub Schedule Reliability [+/- 6 hrs | +/- 24 hrs]
    	image_path_7 = "chart_7.png"  #Hub Productivity [VNMPH*]
    	image_path_8 = "chart_8.png"  #Hub TS Dwell Time [days]
    	image_path_9 = "fotter_1.png" #Share Of Voyages Started On Time fotter
    	image_path_10 = "fotter_2.png" #All Gemini Port Calls | HL vs. Partner | Proforma SR [+/- 1 Day] fotter
    	image_path_11 = "fotter_3.png" #Global | Proforma SR [+/- 1 Day] fotter
    	image_path_12 = "fotter_4.png" #No. Of Gemini Port Calls | HL vs. Partner fotter
    	image_path_13 = "fotter_5.png" #Gemini By Trade | Proforma Schedule Reliability [+/- 1 Day] for Mainliner and Shuttles fotter
    	image_path_14 = "fotter_6.png" #Hub Schedule Reliability [+/- 6 hrs | +/- 24 hrs] fotter
    	image_path_15 = "fotter_7.png" #Hub Productivity [VNMPH*] fotter
   	image_path_16 = "fotter_8.png" #Hub TS Dwell Time [days] fotter
    	image_path_17 = "table_1_fotter.png" #Hub Schedule Reliability [+/- 6 hrs | +/- 24 hrs] table fotter
    	image_path_18 = "table_2_fotter.png" #Hub Productivity [VNMPH*] table fotter
    	image_path_19 = "table_3_fotter.png" #Hub TS Dwell Time [days] table fotter

    	image_path_20 =  "port_1.png"
    	image_path_21 =  "port_2.png"
    	##Tables
    	image_path_22 = "table_1_1.png"
    	image_path_23 = "table_1_2.png"
    	image_path_24 = "table_2_1.png"
    	image_path_25 = "table_2_2.png"
    	image_path_26 = "table_3_1.png" #Hub TS Dwell Time [days] table_1
    	image_path_27 = "table_3_2.png"#Hub TS Dwell Time [days] table_2

    

    
    # Add the first image to the 5th slide at a specific position and size
    	slide_4 = ppt.slides[3]
    	slide_5 = ppt.slides[4]
    	slide_6 = ppt.slides[5]
    	slide_7 = ppt.slides[6]
    	slide_8 = ppt.slides[7]
    	slide_9 = ppt.slides[8]

    ##First Slide
    	add_image_to_slide(slide_4, image_path_20, left_cm=1.76, top_cm=12.12,width_cm=10.04, height_cm=0.73)
    	add_image_to_slide(slide_4, image_path_21, left_cm=14.45, top_cm=12.21,width_cm=10.01, height_cm=0.64)
    	add_image_to_slide(slide_4, image_path_1, left_cm=0.91, top_cm=7.14, width_cm=10.85, height_cm=5.2)
    	add_image_to_slide(slide_4, image_path_2, left_cm=13.39, top_cm=6.33, width_cm=11, height_cm=6.03)
    	add_image_to_slide(slide_4, image_path_9, left_cm=0.91, top_cm=12.66, width_cm=10.89, height_cm=0.38)
    	add_image_to_slide(slide_4, image_path_10, left_cm=13.57, top_cm=12.66,width_cm=10.89, height_cm=0.38)



    ##Second Slide
    	add_image_to_slide(slide_5, image_path_3, left_cm=0.96, top_cm=6.41, width_cm=12.22, height_cm=5.97)
    	add_image_to_slide(slide_5, image_path_4, left_cm=13.39, top_cm=6.42, width_cm=10.95, height_cm=5.99)
    	add_image_to_slide(slide_5, image_path_12, left_cm=0.88, top_cm=12.53, width_cm=11.01, height_cm=0.32)
    	add_image_to_slide(slide_5, image_path_12, left_cm=13.36, top_cm=12.53, width_cm=11.01, height_cm=0.32)


    ##Third SLide
    	add_image_to_slide(slide_6, image_path_5, left_cm=0.92, top_cm=2.18, width_cm=17.11, height_cm=10.79)
    	add_rectangle_to_slide(slide_6, 6.65, 9.3, 11.82, 3.54)
    	##add_image_to_slide(slide_6, image_path_13, left_cm=0.92, top_cm=12.49, width_cm=17.67, height_cm=0.46)
    
    ## Fourth Slide
    	add_image_to_slide(slide_7, image_path_6, left_cm=1.04, top_cm=6.15, width_cm=8.32, height_cm=6.04)
    	add_image_to_slide(slide_7, image_path_14, left_cm=1.04, top_cm=12.29, width_cm=8.35, height_cm=0.64)
    	add_image_to_slide(slide_7, image_path_17, left_cm=10.19, top_cm=12.62, width_cm=12.62, height_cm=0.3)
    	add_image_to_slide(slide_7, image_path_23, left_cm=10.22, top_cm=6.35, width_cm=12.38, height_cm=6.18)
    	add_image_to_slide(slide_7, image_path_22, left_cm=10.22, top_cm=3.8, width_cm=12.38, height_cm=6.18)
    	add_rectangle_to_slide(slide_7, 22.33, 2.97, 1, 9.8)
    	add_rectangle_to_slide(slide_7, 10.22, 3.69, 6.76, 0.42)
    ##Add table

    ## Fifth Slide
    	add_image_to_slide(slide_8, image_path_7, left_cm=0.98, top_cm=6.14, width_cm=8.38, height_cm=6.11)
    	add_image_to_slide(slide_8, image_path_15, left_cm=1, top_cm=12.36, width_cm=8.38, height_cm=0.61)
    	add_image_to_slide(slide_8, image_path_18, left_cm=10.17, top_cm=12.38, width_cm=12.8, height_cm=0.61)
    	add_image_to_slide(slide_8, image_path_25, left_cm=10.15, top_cm=5.55, width_cm=13.4, height_cm=6.68)
    	add_image_to_slide(slide_8, image_path_24, left_cm=10.15, top_cm=3.66, width_cm=13.4, height_cm=7.03)
    	add_rectangle_to_slide(slide_8, 22.77, 4.15 , 0.97, 8.22)
    	add_rectangle_to_slide(slide_8, 10.17, 3.48 , 4.52, 0.62)
    ##Add table

    ## sixth Slide
    	add_image_to_slide(slide_9, image_path_8, left_cm=1.02, top_cm=6.13, width_cm=7.8, height_cm=6.45)
    	add_image_to_slide(slide_9, image_path_16, left_cm=1.02, top_cm=12.69, width_cm=7.8, height_cm=0.26)
    	add_image_to_slide(slide_9, image_path_16, left_cm=9.56, top_cm=12.71, width_cm=7.8, height_cm=0.26)
    	add_image_to_slide(slide_9, image_path_27, left_cm=9.55, top_cm=5.77, width_cm=11.55, height_cm=6.77)
    	add_image_to_slide(slide_9, image_path_26, left_cm=9.55, top_cm=3.5, width_cm=11.56, height_cm=6.77)
    	add_rectangle_to_slide(slide_9, 21.12, 4.13, 8.58, 0.65)
    	add_rectangle_to_slide(slide_9, 9.52, 3.44, 3.95, 0.37)

        output_filename = "Gemini Ops Transition Reporting _ CW.pptx"
        prs.save(output_filename)

        with open(output_filename, "rb") as f:
            st.success("PPT generated successfully. Click below to download.")
            st.download_button("Download PPT", f, file_name=output_filename)

    except Exception as e:
        st.error(f"An error occurred: {e}")

    finally:
        driver.quit()
        driver_2.quit()
        st.info("Browser sessions closed.")
